"""
SHAMEL Graduation Presentation — python-pptx builder
Formal dark-navy + gold color scheme, comprehensive problem→solution mapping
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as pns
from lxml import etree
import copy

# ── Colors ──────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0B, 0x14, 0x2B)   # near-black navy bg
NAVY2       = RGBColor(0x14, 0x24, 0x45)   # card bg
GOLD        = RGBColor(0xC9, 0xA8, 0x2C)   # primary accent
GOLD2       = RGBColor(0xE8, 0xCA, 0x6A)   # lighter gold
WHITE       = RGBColor(0xF4, 0xF6, 0xFB)   # body text
MUTED       = RGBColor(0x8A, 0x9B, 0xBF)   # secondary text
EMERALD     = RGBColor(0x10, 0xB9, 0x81)   # success / solution
RED_LIGHT   = RGBColor(0xFC, 0xA5, 0xA5)   # problem highlight
BLUE_LIGHT  = RGBColor(0x60, 0xA5, 0xFA)   # info

W = Inches(13.33)
H = Inches(7.5)

# ── Helpers ─────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_w=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE.RECTANGLE
        l, t, w, h
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "Calibri"
    return txb

def add_gradient_line(slide, l, t, w, h=Inches(0.04)):
    """Thin gold accent line."""
    shape = add_rect(slide, l, t, w, h, fill_color=GOLD)
    return shape

def slide_header(slide, tag, title, subtitle=None):
    """Consistent header block for content slides."""
    # tag pill
    tag_box = add_rect(slide, Inches(.5), Inches(.35), Inches(2.2), Inches(.32),
                        fill_color=RGBColor(0x1E, 0x30, 0x55))
    add_text(slide, tag, Inches(.5), Inches(.33), Inches(2.2), Inches(.35),
             size=Pt(9), bold=True, color=GOLD2, align=PP_ALIGN.CENTER)
    # title
    add_text(slide, title, Inches(.5), Inches(.75), Inches(12.3), Inches(.9),
             size=Pt(34), bold=True, color=WHITE)
    # gold line
    add_gradient_line(slide, Inches(.5), Inches(1.65), Inches(1.2))
    if subtitle:
        add_text(slide, subtitle, Inches(.5), Inches(1.72), Inches(11), Inches(.45),
                 size=Pt(15), color=MUTED)

# ── SLIDE BUILDERS ───────────────────────────────────────────────────────────

def s01_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    # Left accent bar
    add_rect(sl, Inches(0), Inches(0), Inches(.12), H, fill_color=GOLD)

    # Background subtle rect
    add_rect(sl, Inches(.12), Inches(0), Inches(13.21), H,
             fill_color=NAVY)

    # Large project name
    add_text(sl, "SHAMEL", Inches(.6), Inches(1.2), Inches(8), Inches(1.8),
             size=Pt(88), bold=True, color=GOLD)

    # Full name
    add_text(sl, "Smart Holistic Attendance\nManagement & Compliance System",
             Inches(.6), Inches(2.9), Inches(9), Inches(1.4),
             size=Pt(22), bold=False, color=WHITE)

    add_gradient_line(sl, Inches(.6), Inches(4.38), Inches(5))

    # Team info block
    add_text(sl, "Team Members", Inches(.6), Inches(4.55), Inches(3), Inches(.3),
             size=Pt(10), bold=True, color=GOLD2)
    add_text(sl, "Member 1  ·  Member 2  ·  Member 3  ·  Member 4",
             Inches(.6), Inches(4.85), Inches(7), Inches(.3),
             size=Pt(13), color=WHITE)

    add_text(sl, "Supervisor", Inches(.6), Inches(5.25), Inches(3), Inches(.3),
             size=Pt(10), bold=True, color=GOLD2)
    add_text(sl, "Dr. Supervisor Name",
             Inches(.6), Inches(5.55), Inches(5), Inches(.3),
             size=Pt(13), color=WHITE)

    add_text(sl, "College of Computer Science & IT  ·  2024 / 2025",
             Inches(.6), Inches(6.7), Inches(10), Inches(.3),
             size=Pt(11), color=MUTED)

    # Right side: geometric decoration
    add_rect(sl, Inches(9.8), Inches(1.0), Inches(3.0), Inches(5.0),
             fill_color=RGBColor(0x12, 0x22, 0x40))
    add_text(sl, "🎓", Inches(9.8), Inches(2.5), Inches(3.0), Inches(2.0),
             size=Pt(90), align=PP_ALIGN.CENTER, color=GOLD)

def s02_overview(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    add_rect(sl, Inches(0), Inches(0), Inches(.08), H, fill_color=GOLD)
    slide_header(sl, "OVERVIEW", "What SHAMEL Really Is",
                 "A unified system that digitalises every layer of university operations")

    cards = [
        ("🎯", "Biometric\nAccess Control",   "Gate + Classroom"),
        ("📊", "Academic\nManagement",         "Grades · Excuses · Reports"),
        ("🔄", "Real-Time\nMonitoring",        "WebSocket live dashboards"),
        ("📱", "Multi-Platform",               "Web PWA + Android App"),
        ("🔒", "Role-Based\nSecurity",         "5 isolated role scopes"),
        ("📡", "Offline-First",                "Works without internet"),
    ]
    cw = Inches(1.9); ch = Inches(1.75)
    cols = 3
    for i, (icon, title, sub) in enumerate(cards):
        col = i % cols; row = i // cols
        l = Inches(.5) + col * (cw + Inches(.18))
        t = Inches(2.1) + row * (ch + Inches(.15))
        add_rect(sl, l, t, cw, ch, fill_color=NAVY2)
        add_text(sl, icon, l, t + Inches(.1), cw, Inches(.55),
                 size=Pt(26), align=PP_ALIGN.CENTER, color=WHITE)
        add_text(sl, title, l, t + Inches(.62), cw, Inches(.65),
                 size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, sub, l, t + Inches(1.25), cw, Inches(.35),
                 size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)

def make_problem_slide(prs, tag, title, problems, solutions, icon="⚙️"):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    add_rect(sl, Inches(0), Inches(0), Inches(.08), H, fill_color=GOLD)
    slide_header(sl, tag, title)

    # Two-column: Problems (left) | Solutions (right)
    col_w = Inches(5.8)
    gap   = Inches(.25)
    lx    = Inches(.5)
    rx    = lx + col_w + gap + Inches(.15)
    ty    = Inches(2.0)

    # Column headers
    add_rect(sl, lx, ty, col_w, Inches(.35), fill_color=RGBColor(0x3B, 0x1A, 0x1A))
    add_text(sl, "❌  The Problem", lx + Inches(.1), ty + Inches(.04), col_w, Inches(.3),
             size=Pt(11), bold=True, color=RED_LIGHT)

    add_rect(sl, rx, ty, col_w, Inches(.35), fill_color=RGBColor(0x0E, 0x2E, 0x23))
    add_text(sl, "✅  SHAMEL's Solution", rx + Inches(.1), ty + Inches(.04), col_w, Inches(.3),
             size=Pt(11), bold=True, color=EMERALD)

    row_h = Inches(.72)
    for i, (prob, sol) in enumerate(zip(problems, solutions)):
        y = ty + Inches(.42) + i * row_h
        bg = RGBColor(0x1A, 0x10, 0x10) if i % 2 == 0 else RGBColor(0x16, 0x0E, 0x0E)
        add_rect(sl, lx, y, col_w, row_h - Inches(.06), fill_color=bg)
        add_text(sl, f"• {prob}", lx + Inches(.12), y + Inches(.08),
                 col_w - Inches(.15), row_h - Inches(.12),
                 size=Pt(13), color=WHITE, wrap=True)

        bg2 = RGBColor(0x0A, 0x1E, 0x18) if i % 2 == 0 else RGBColor(0x0C, 0x1C, 0x16)
        add_rect(sl, rx, y, col_w, row_h - Inches(.06), fill_color=bg2)
        add_text(sl, f"→ {sol}", rx + Inches(.12), y + Inches(.08),
                 col_w - Inches(.15), row_h - Inches(.12),
                 size=Pt(13), color=EMERALD, wrap=True)

def s03_identity(prs):
    make_problem_slide(prs,
        "PROBLEM → SOLUTION",
        "Identity & Attendance Fraud",
        [
            "Proxy attendance — anyone signs for another",
            "No biometric verification at entry",
            "Forged paper registers impossible to detect",
            "No cooldown — same person re-enters repeatedly",
        ],[
            "InsightFace buffalo_s — 512-dim cosine, spoofing fails",
            "Two-phase gate: MediaPipe (4ms) + InsightFace (9ms)",
            "Every mark tied to a face embedding, not a signature",
            "30-second server-side cooldown per person per session",
        ]
    )

def s04_visibility(prs):
    make_problem_slide(prs,
        "PROBLEM → SOLUTION",
        "Zero Real-Time Visibility",
        [
            "Coordinators learn about absences days later",
            "No live view of who is in which classroom",
            "Gate entry logs checked manually at end of day",
            "Students unaware of their own attendance status",
        ],[
            "Django Channels WebSocket — events pushed instantly",
            "Live dashboard: per-classroom headcount, live chart",
            "Gate logs update in real-time, filterable by role",
            "Student portal shows % and schedule gaps immediately",
        ]
    )

def s05_records(prs):
    make_problem_slide(prs,
        "PROBLEM → SOLUTION",
        "Paper Records & Reporting",
        [
            "Attendance sheets lost, damaged, or tampered",
            "Reports take days to compile manually",
            "No analytics on at-risk students",
            "Audit trail nonexistent when disputes arise",
        ],[
            "All records in PostgreSQL with full timestamp chain",
            "One-click PDF (WeasyPrint) · Excel · CSV exports",
            "Coordinator dashboard flags students below threshold",
            "AuditLog model records every sensitive action + actor",
        ]
    )

def s06_access(prs):
    make_problem_slide(prs,
        "PROBLEM → SOLUTION",
        "Access Control & Data Silos",
        [
            "Staff see data they shouldn't (privacy violations)",
            "No separation between administrative & academic roles",
            "Gate operators had access to grade records",
            "Coordinators were seeing other colleges' data",
        ],[
            "5-role RBAC: Admin · Coordinator · Teacher · Student · Gate",
            "Coordinator scoped to one College only — enforced in DB",
            "Gate role: access logs only, zero academic data",
            "Every view decorated with role-check, API same",
        ]
    )

def s07_excuse(prs):
    make_problem_slide(prs,
        "PROBLEM → SOLUTION",
        "Medical Excuse & Communication Gaps",
        [
            "Paper excuse forms lost or delayed in bureaucracy",
            "No document trail for medical absence disputes",
            "Teachers unaware of approved excuses when grading",
            "No automated notification when student hits limit",
        ],[
            "Digital excuse portal: upload docs, approve/reject online",
            "Every excuse stored with file, status, timestamps",
            "Approved excuses auto-visible in teacher's record view",
            "SMTP email alert sent when attendance falls below 75%",
        ]
    )

def s08_offline(prs):
    make_problem_slide(prs,
        "PROBLEM → SOLUTION",
        "Connectivity & Offline Reliability",
        [
            "Single internet outage breaks entire attendance flow",
            "VPS unreachable = system completely down",
            "Mobile networks in lecture halls unreliable",
            "Data loss if session ends before upload completes",
        ],[
            "IndexedDB offline queue on web — auto-sync on reconnect",
            "Auto-detect PostgreSQL → SQLite fallback at startup",
            "Flutter: sqflite local DB + connectivity_plus sync queue",
            "Idempotent sync API — duplicate submissions are safe",
        ]
    )

def s09_grades(prs):
    make_problem_slide(prs,
        "PROBLEM → SOLUTION",
        "Grade & Eligibility Management",
        [
            "Manual eligibility calculation error-prone",
            "Students surprised by exam bans at the last moment",
            "No link between attendance record and final grade",
            "Exam seat assignment done manually on paper",
        ],[
            "Attendance % auto-calculated per course per student",
            "Ineligibility flag raised early — email sent immediately",
            "Grade model linked to attendance & excuse records",
            "ExamSeat model: auto-assign + PDF seating chart export",
        ]
    )

def s10_arch(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    add_rect(sl, Inches(0), Inches(0), Inches(.08), H, fill_color=GOLD)
    slide_header(sl, "ARCHITECTURE", "How the System Connects",
                 "Every component designed for resilience and scale")

    nodes = [
        (Inches(.5),  Inches(2.2), "🌐 Web Browser\n(PWA)",        NAVY2,   BLUE_LIGHT),
        (Inches(.5),  Inches(3.8), "📱 Flutter App\n(Android)",    NAVY2,   BLUE_LIGHT),
        (Inches(.5),  Inches(5.4), "🎥 Gate Camera\n(RTSP)",       NAVY2,   BLUE_LIGHT),
        (Inches(3.4), Inches(3.4), "🔀 Nginx\n+ Daphne :9000",    RGBColor(0x12,0x22,0x40), GOLD),
        (Inches(6.2), Inches(3.4), "🌿 Django\nAPI + Views",       RGBColor(0x10,0x26,0x1A), EMERALD),
        (Inches(9.0), Inches(2.2), "🧠 InsightFace\nONNX Engine",  RGBColor(0x0E,0x28,0x20), EMERALD),
        (Inches(9.0), Inches(3.8), "🐘 PostgreSQL\n+ SQLite",      RGBColor(0x0E,0x28,0x20), EMERALD),
        (Inches(9.0), Inches(5.4), "⚡ Redis\nChannels WS",         RGBColor(0x0E,0x28,0x20), EMERALD),
    ]
    bw = Inches(2.5); bh = Inches(1.0)
    for (l, t, txt, bg, col) in nodes:
        add_rect(sl, l, t, bw, bh, fill_color=bg)
        add_text(sl, txt, l + Inches(.08), t + Inches(.08),
                 bw - Inches(.1), bh - Inches(.1),
                 size=Pt(12), bold=True, color=col, align=PP_ALIGN.CENTER, wrap=True)

    # Arrows (simple text)
    for y in [Inches(2.68), Inches(3.93), Inches(5.86)]:
        add_text(sl, "→", Inches(3.05), y, Inches(.4), Inches(.3),
                 size=Pt(18), color=MUTED, align=PP_ALIGN.CENTER)
    for y in [Inches(2.68), Inches(3.93), Inches(5.86)]:
        add_text(sl, "→", Inches(5.85), y, Inches(.4), Inches(.3),
                 size=Pt(18), color=MUTED, align=PP_ALIGN.CENTER)
    for y in [Inches(2.68), Inches(3.93), Inches(5.86)]:
        add_text(sl, "→", Inches(8.65), y, Inches(.4), Inches(.3),
                 size=Pt(18), color=MUTED, align=PP_ALIGN.CENTER)

def s11_tech(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    add_rect(sl, Inches(0), Inches(0), Inches(.08), H, fill_color=GOLD)
    slide_header(sl, "TECH STACK", "Built with the Right Tools")

    stack = [
        ("🐍", "Python 3.12",      "Runtime"),
        ("🌿", "Django 4.x",       "Web + API + WS"),
        ("📱", "Flutter 3.41",     "Android App"),
        ("🧠", "InsightFace ONNX", "Face Recognition"),
        ("👁️", "MediaPipe",        "Client Detection"),
        ("🐘", "PostgreSQL",       "Primary DB"),
        ("💾", "SQLite",           "Offline Fallback"),
        ("⚡", "Redis + Channels", "WebSocket Broker"),
        ("🐳", "Docker",           "Containerisation"),
        ("🔀", "Nginx + Daphne",   "ASGI Reverse Proxy"),
        ("📡", "PWA / SW",         "Service Worker"),
        ("🔑", "Bearer Token",     "API Auth"),
    ]
    cw = Inches(2.8); ch = Inches(.92)
    cols = 4
    for i, (icon, name, role) in enumerate(stack):
        col = i % cols; row = i // cols
        l = Inches(.45) + col * (cw + Inches(.12))
        t = Inches(2.1) + row * (ch + Inches(.12))
        add_rect(sl, l, t, cw, ch, fill_color=NAVY2)
        add_text(sl, icon + "  " + name, l + Inches(.12), t + Inches(.1),
                 cw - Inches(.15), Inches(.42),
                 size=Pt(14), bold=True, color=WHITE)
        add_text(sl, role, l + Inches(.12), t + Inches(.52),
                 cw - Inches(.15), Inches(.32),
                 size=Pt(11), color=MUTED)

def s12_perf(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    add_rect(sl, Inches(0), Inches(0), Inches(.08), H, fill_color=GOLD)
    slide_header(sl, "RESULTS", "System Performance")

    stats = [
        ("≥ 95%",  "Face Match\nAccuracy",   GOLD),
        ("~4 ms",  "Face Detect\n(InsightFace)", EMERALD),
        ("~9 ms",  "Embedding\nEncode",       EMERALD),
        ("13 ms",  "End-to-End\nInference",   BLUE_LIGHT),
        ("5",      "Isolated\nRole Scopes",   GOLD),
        ("100%",   "Offline-First\nCapable",  EMERALD),
    ]
    sw = Inches(1.9); sh = Inches(1.9)
    cols = 3
    for i, (num, label, col) in enumerate(stats):
        c = i % cols; r = i // cols
        l = Inches(.5) + c * (sw + Inches(.25))
        t = Inches(2.1) + r * (sh + Inches(.15))
        add_rect(sl, l, t, sw, sh, fill_color=NAVY2)
        add_text(sl, num, l, t + Inches(.25), sw, Inches(.85),
                 size=Pt(36), bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(sl, label, l, t + Inches(1.1), sw, Inches(.65),
                 size=Pt(12), color=MUTED, align=PP_ALIGN.CENTER)

def s13_future(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    add_rect(sl, Inches(0), Inches(0), Inches(.08), H, fill_color=GOLD)
    slide_header(sl, "FUTURE SCOPE", "Where SHAMEL Goes Next")

    items = [
        ("🔍", "pgvector Integration",    "GPU-accelerated cosine search — O(1) at any scale"),
        ("📝", "Exam Proctoring",          "Continuous face verification + auto seat assignment"),
        ("💬", "SMS / WhatsApp Alerts",    "Instant parent notification on absence breach"),
        ("🏛️", "Multi-Campus Federation", "Shared embedding store across university branches"),
    ]
    for i, (icon, title, body) in enumerate(items):
        y = Inches(2.1) + i * Inches(1.12)
        add_rect(sl, Inches(.5), y, Inches(12.3), Inches(1.0), fill_color=NAVY2)
        add_text(sl, icon, Inches(.65), y + Inches(.18), Inches(.6), Inches(.65),
                 size=Pt(26), color=GOLD)
        add_text(sl, title, Inches(1.35), y + Inches(.1), Inches(4), Inches(.42),
                 size=Pt(16), bold=True, color=WHITE)
        add_text(sl, body,  Inches(1.35), y + Inches(.52), Inches(10.8), Inches(.4),
                 size=Pt(13), color=MUTED)

def s14_thanks(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    add_rect(sl, Inches(0), Inches(0), Inches(.08), H, fill_color=GOLD)

    add_rect(sl, Inches(.12), Inches(0), Inches(13.21), H, fill_color=NAVY)
    add_text(sl, "SHAMEL", Inches(.5), Inches(1.5), Inches(12.3), Inches(1.8),
             size=Pt(90), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_gradient_line(sl, Inches(4.5), Inches(3.4), Inches(4.3))
    add_text(sl, "Thank You  ·  Q & A", Inches(.5), Inches(3.6), Inches(12.3), Inches(.7),
             size=Pt(28), bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "shamel.sd  ·  Python · Django · Flutter · InsightFace ONNX",
             Inches(.5), Inches(4.5), Inches(12.3), Inches(.4),
             size=Pt(13), color=MUTED, align=PP_ALIGN.CENTER)
    add_text(sl, "College of Computer Science & IT  ·  2024 / 2025",
             Inches(.5), Inches(6.7), Inches(12.3), Inches(.35),
             size=Pt(11), color=RGBColor(0x3A, 0x4A, 0x65), align=PP_ALIGN.CENTER)

# ── MAIN ────────────────────────────────────────────────────────────────────
prs = new_prs()

s01_title(prs)
s02_overview(prs)
s03_identity(prs)
s04_visibility(prs)
s05_records(prs)
s06_access(prs)
s07_excuse(prs)
s08_offline(prs)
s09_grades(prs)
s10_arch(prs)
s11_tech(prs)
s12_perf(prs)
s13_future(prs)
s14_thanks(prs)

import os, sys
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "SHAMEL_Presentation.pptx")
prs.save(out)
sys.stdout.buffer.write(b"Saved: " + out.encode('utf-8') + b"\n")
